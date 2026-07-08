from __future__ import annotations

import csv
from dataclasses import dataclass
from io import BytesIO, StringIO
from pathlib import Path

from core.settings import get_settings


@dataclass
class ParsedDocument:
    content: str
    metadata: dict
    source_type: str


def parse_text(text: str) -> str:
    return text.strip()


def parse_pdf(file_bytes: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ValueError("PDF ingest is disabled. Install 'pypdf' to enable .pdf parsing.") from exc
    reader = PdfReader(BytesIO(file_bytes))
    pages: list[str] = []
    for page in reader.pages:
        pages.append((page.extract_text() or "").strip())
    return "\n\n".join([page for page in pages if page])


def parse_excel(file_bytes: bytes) -> str:
    try:
        import pandas as pd
    except ImportError as exc:
        raise ValueError("Excel ingest is disabled. Install 'pandas' and 'openpyxl' to enable .xlsx/.xls parsing.") from exc
    sheet_map = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
    parts: list[str] = []
    for sheet_name, dataframe in sheet_map.items():
        as_text = dataframe.fillna("").to_csv(index=False)
        parts.append(f"[Sheet: {sheet_name}]\n{as_text}")
    return "\n\n".join(parts)


def parse_csv(file_bytes: bytes) -> str:
    decoded = file_bytes.decode("utf-8", errors="ignore")
    rows = list(csv.reader(StringIO(decoded)))
    output = StringIO()
    writer = csv.writer(output)
    for row in rows:
        writer.writerow(row)
    return output.getvalue().strip()


def parse_ppt(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValueError("PowerPoint ingest is disabled. Install 'python-pptx' to enable .ppt/.pptx parsing.") from exc
    presentation = Presentation(BytesIO(file_bytes))
    slides: list[str] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                slide_text.append(text.strip())
        if slide_text:
            slides.append(f"[Slide {idx}]\n" + "\n".join(slide_text))
    return "\n\n".join(slides)


def _chunk_text(text: str, chunk_size: int, overlap: int) -> list[str]:
    if not text:
        return []
    if chunk_size <= 0:
        return [text]

    chunks: list[str] = []
    start = 0
    step = max(1, chunk_size - max(0, overlap))
    while start < len(text):
        chunk = text[start : start + chunk_size].strip()
        if chunk:
            chunks.append(chunk)
        start += step
    return chunks


def _parse_by_extension(file_bytes: bytes, extension: str) -> str:
    ext = extension.lower()
    if ext in {".txt", ".md"}:
        return parse_text(file_bytes.decode("utf-8", errors="ignore"))
    if ext == ".pdf":
        return parse_pdf(file_bytes)
    if ext == ".csv":
        return parse_csv(file_bytes)
    if ext in {".xlsx", ".xls"}:
        return parse_excel(file_bytes)
    if ext in {".ppt", ".pptx"}:
        return parse_ppt(file_bytes)
    raise ValueError(f"Unsupported file extension: {extension}")


def parse_file(file_bytes: bytes, filename: str):
    settings = get_settings()
    extension = Path(filename).suffix
    full_text = _parse_by_extension(file_bytes, extension)
    source_type = "consultant_profile"
    chunks = _chunk_text(
        text=full_text,
        chunk_size=settings.ingest_chunk_size,
        overlap=settings.ingest_chunk_overlap,
    )

    for idx, chunk in enumerate(chunks):
        yield ParsedDocument(
            content=chunk,
            metadata={
                "source": filename,
                "source_type": source_type,
                "chunk_index": idx,
                "chunk_count": len(chunks),
            },
            source_type=source_type,
        )
